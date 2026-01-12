"""
PowerPoint Presentation Generator for AI-Powered Wireless Firewall
Creates a professional presentation for the Wireless Technologies course
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Alias for convenience
RgbColor = RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(26, 26, 46)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(100, 200, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, emoji=""):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(26, 26, 46)
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(45, 45, 80)
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Bullet points
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = RgbColor(220, 220, 220)
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, emoji=""):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(26, 26, 46)
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(45, 45, 80)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(100, 200, 255)
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.2), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(220, 220, 220)
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 150, 100)
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.2), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RgbColor(220, 220, 220)
        p.space_after = Pt(8)
    
    return slide

def add_stats_slide(prs, title, stats, emoji=""):
    """Add a slide with statistics boxes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(26, 26, 46)
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(45, 45, 80)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{emoji} {title}" if emoji else title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Stats boxes
    colors = [
        RgbColor(46, 125, 50),   # Green
        RgbColor(21, 101, 192),  # Blue
        RgbColor(211, 47, 47),   # Red
        RgbColor(255, 143, 0),   # Orange
    ]
    
    box_width = Inches(2.1)
    box_height = Inches(1.8)
    start_x = Inches(0.5)
    gap = Inches(0.2)
    
    for i, (label, value) in enumerate(stats[:4]):
        x = start_x + i * (box_width + gap)
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()
        
        # Value
        val_box = slide.shapes.add_textbox(x, Inches(2), box_width, Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        lbl_box = slide.shapes.add_textbox(x, Inches(2.8), box_width, Inches(0.6))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(26, 26, 46)
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(45, 45, 80)
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    # Docker container boxes
    containers = [
        ("NGINX Dashboard", ":80", Inches(0.5), RgbColor(76, 175, 80)),
        ("AI-Engine (FastAPI)", ":8000", Inches(3.4), RgbColor(33, 150, 243)),
        ("Redis Cache", ":6379", Inches(6.3), RgbColor(255, 87, 34)),
    ]
    
    for name, port, x, color in containers:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), Inches(2.6), Inches(1.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        # Name
        name_box = slide.shapes.add_textbox(x, Inches(2.3), Inches(2.6), Inches(0.6))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Port
        port_box = slide.shapes.add_textbox(x, Inches(2.9), Inches(2.6), Inches(0.4))
        tf = port_box.text_frame
        p = tf.paragraphs[0]
        p.text = port
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    # ML Models box
    ml_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4), Inches(5), Inches(1.8)
    )
    ml_box.fill.solid()
    ml_box.fill.fore_color.rgb = RgbColor(103, 58, 183)
    ml_box.line.fill.background()
    
    ml_title = slide.shapes.add_textbox(Inches(2.5), Inches(4.2), Inches(5), Inches(0.5))
    tf = ml_title.text_frame
    p = tf.paragraphs[0]
    p.text = "ML Ensemble Model"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    ml_content = slide.shapes.add_textbox(Inches(2.5), Inches(4.7), Inches(5), Inches(1))
    tf = ml_content.text_frame
    p = tf.paragraphs[0]
    p.text = "XGBoost (70%) + Isolation Forest (30%)"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(220, 220, 220)
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "84 Features | CIC-IDS2017 Dataset"
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER
    
    # Network label
    net_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.4))
    tf = net_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Docker Network: firewall-net (172.28.0.0/16)"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(150, 150, 150)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Create the full presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "AI-Powered Wireless Firewall",
        "Wireless Technologies Project 2026"
    )
    
    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "AI-based firewall that analyzes network traffic",
        "Automatically distinguishes between normal and malicious traffic",
        "Uses Machine Learning to detect attack patterns",
        "Real-time dashboard for monitoring",
        "100% attack detection rate on trained attack types"
    ], "🎯")
    
    # Slide 3: Problem Statement
    add_two_column_slide(
        prs, "Problem Statement",
        "Wireless Threats",
        [
            "WiFi Deauth Attacks",
            "Evil Twin Access Points",
            "KRACK WPA2 Attacks",
            "Bluetooth Hijacking",
            "IoT Device Attacks"
        ],
        "The Challenge",
        [
            "Traditional firewalls use fixed rules",
            "Can't detect unknown attacks",
            "Manual configuration required",
            "No real-time adaptation",
            "High false positive rates"
        ],
        "⚠️"
    )
    
    # Slide 4: Solution
    add_content_slide(prs, "Our Solution", [
        "Machine Learning analyzes 84 network features per flow",
        "Ensemble model combines XGBoost + Isolation Forest",
        "Trained on 2.8 million network flows (CIC-IDS2017)",
        "Real-time predictions in milliseconds",
        "Web dashboard for live monitoring",
        "Docker containerized for easy deployment"
    ], "💡")
    
    # Slide 5: Architecture
    add_architecture_slide(prs)
    
    # Slide 6: ML Model Details
    add_two_column_slide(
        prs, "Machine Learning Model",
        "XGBoost Classifier (70%)",
        [
            "Gradient Boosted Decision Trees",
            "100 trees, max depth 6",
            "Learns complex attack patterns",
            "High accuracy on known attacks",
            "Output: Probability 0-100%"
        ],
        "Isolation Forest (30%)",
        [
            "Anomaly Detection Algorithm",
            "100 isolation trees",
            "Detects unknown attacks",
            "No labels required",
            "Output: Anomaly score"
        ],
        "🤖"
    )
    
    # Slide 7: Dataset
    add_stats_slide(prs, "Training Dataset: CIC-IDS2017", [
        ("Total Flows", "2.8M"),
        ("Features", "84"),
        ("Attack Types", "14"),
        ("File Size", "1.2GB"),
    ], "📊")
    
    # Slide 8: Wireless Attacks Detected
    add_two_column_slide(
        prs, "Wireless Attack Detection",
        "WiFi Attacks (100%)",
        [
            "WiFi Deauth Attack",
            "Evil Twin AP",
            "KRACK Attack",
            "WiFi Jamming",
            "PMKID Attack",
            "Wardriving Probe"
        ],
        "Bluetooth/IoT (100%)",
        [
            "Bluetooth Hijacking",
            "IoT Zigbee Attack",
            "Smart Home Attacks",
            "",
            "✅ All detected with >90% confidence"
        ],
        "📡"
    )
    
    # Slide 9: Dashboard
    add_content_slide(prs, "Real-time Dashboard", [
        "Statistics cards: Total flows, Benign, Malicious, Attack Types",
        "Flow Classification doughnut chart",
        "Threat Timeline line chart",
        "Recent Alerts panel with attack details",
        "Auto-refresh every 1 second via HTTP polling",
        "Dark theme responsive design"
    ], "📊")
    
    # Slide 10: Results
    add_stats_slide(prs, "Results", [
        ("Attack Detection", "100%"),
        ("Benign Accuracy", "100%"),
        ("False Positives", "0%"),
        ("Response Time", "<50ms"),
    ], "📈")
    
    # Slide 11: Demo
    add_content_slide(prs, "Live Demo", [
        "1. Start containers: docker-compose up -d",
        "2. Open dashboard: http://localhost:80",
        "3. Run attack simulator: python test_wireless_attacks.py",
        "4. Watch real-time detections on dashboard",
        "",
        "8 attack types × 3 samples = 24 attacks detected",
        "5 benign types × 3 samples = 15 benign flows allowed"
    ], "🎬")
    
    # Slide 12: Future Improvements
    add_content_slide(prs, "Future Improvements", [
        "Online learning - Model adapts to new attacks",
        "Mobile app for smartphone monitoring",
        "Integration with existing network hardware",
        "Cloud deployment as SaaS solution",
        "Extended IoT protocol support"
    ], "🔮")
    
    # Slide 13: Conclusion
    add_content_slide(prs, "Conclusion", [
        "✅ AI-based firewall analyzing network traffic",
        "✅ 100% detection of wireless attacks",
        "✅ 0% false positives on normal traffic",
        "✅ Real-time dashboard for monitoring",
        "✅ Scalable Docker architecture",
        "",
        "Machine Learning provides powerful network security automation"
    ], "🎓")
    
    # Slide 14: Questions
    add_title_slide(
        prs,
        "Questions?",
        "GitHub: Darkknin12/wireless-AI-firewall"
    )
    
    # Save
    output_path = "AI_Firewall_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
